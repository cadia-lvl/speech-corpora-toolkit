from pptx import Presentation
import pdfplumber


class DocumentReader:
    def text_from_presentation(self, filepath):
        prs = Presentation(filepath)

        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Only add strings with text
                        if run.text:
                            text.append(run.text)
        return text

    def text_from_pdf(self, filepath):
        text = []
        pdf = pdfplumber.open(filepath)
        for page in pdf.pages:
            text.append(page.extract_text())
        return text

    def read(self, filepath) -> list:
        if ".pptx" in filepath:
            return self.text_from_presentation(filepath)
        if ".pdf" in filepath:
            return self.text_from_pdf(filepath)


def main():
    file = "test_pdf.pdf"
    reader = DocumentReader()
    text = reader.read(file)
    for t in text:
        print(t)


if __name__ == "__main__":
    main()
