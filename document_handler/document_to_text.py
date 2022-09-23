#!/usr/bin/env python3

"""
    This class handles converting documents into a txt format.

    For a proof of concept this is tried with the same google presentation
    saved as .pdf and .pptx

    Functionality:
    reader = DocumentReader()
    reader.read(input_filepath)
    # or trim, trims away lines containing words from trim list
    reader.read_and_trim(input_filepath, trim_list)
    reader.save(output_filepath) # -> saves as plain text

    TODO: Clean .pdf output from itemized and numbered list symbols.

"""

___author___ = "Staffan Hedström"
___license___ = "Apache 2.0"
___copyright___ = "2022 Staffan Hedström Reykjavík University"

from pptx import Presentation
import pdfplumber
from docx import Document
import re
from pathlib import Path

# Trim commands, using input1___input2+++command+++
REMOVE_TABBED = "tabs"
REMOVE_WORD = "word"
REMOVE_FROM_TO = "remove_from_to"
CONTAINS_MULTIPLE = (
    "contains_multiple"  # " / ___3" => removes lines containing 3x " / "
)
REMOVE_NUMBER = "remove_number"  # "part ___3" -> removes line where "part " is followed by 3 or more numbers
TIME_TO_TIME = "time_to_time"  # removes 0.00 – 1.34 and similar
EXACT = "exact"
SONGS = "songs"
TO_GUEST = "to_guest"
ALL_UPPERCASE = "all_uppercase"
REMOVE_START_TIME = "start_time"
REMOVE_SPEAKER_IDENTIFIER = "remove_speaker"
ONLY_NUMBERS = "only_numbers"


class DocumentReader:
    def __init__(self) -> None:
        self.text = []
        self.guest = []

    def text_from_presentation(self, filepath) -> list:
        prs = Presentation(filepath)

        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Only add strings with text
                        if run.text:
                            text.append(run.text)
        self.text = text
        return text

    def text_from_pdf(self, filepath) -> list:
        text = []
        pdf = pdfplumber.open(filepath)
        for page in pdf.pages:
            line = page.extract_text()
            # Make each line an item in the text list
            lines = line.split("\n")
            text.extend(lines)
        self.text = text
        return text

    def text_from_docx(self, filepath: str) -> list:
        text = []
        docx = Document(filepath)
        for p in docx.paragraphs:
            line = p.text
            if line == "" and len(text) >= 1 and text[-1] == "":  # Skip empty lines
                continue
            # Perserve tabs
            try:
                indent = p.paragraph_format.left_indent
                if indent:
                    line = "\t" + line
            except ValueError:
                # Value error means something is wrong with python-docx package
                # It did not convert a value properly but an indentation exists
                line = "\t" + line
            text.append(line)  # fixes a few whitespace issues
        self.text = text
        return text

    def trim(self, trim_list) -> str:
        """
        Removes (trims) away lines containing words in the trim list.
        An item in the trim_list can also be a command which will
        be handled seperately.

        Updates the text and returns the trimmed text
        """
        trimmed_text = []

        for line in self.text:
            if line == "" and len(trimmed_text) >= 1 and trimmed_text[-1] != "":
                trimmed_text.append(line)
                continue
            trimmed = self.trim_line(line, trim_list)
            # Skip empty lines
            if trimmed == "":
                continue
            trimmed_text.append(trimmed)

        for trim_item in trim_list:
            command = self.is_trim_command(trim_item)
            if not command == REMOVE_FROM_TO:
                continue
            start, end = trim_item.replace(f"+++{command}+++", "").split("___")

        # Lastly remove empty lines
        trimmed_text = [line for line in trimmed_text if line.strip() != ""]
        self.guest = [line for line in self.guest if line.strip() != ""]
        self.text = trimmed_text
        return trimmed_text

    def trim_line(self, line, trim_list) -> str:
        """
        Trims a single line. Returns an empty string if trimmed.
        """
        for trim_item in trim_list:
            command = self.is_trim_command(trim_item)
            second_command = None
            if len(command.split(",")) == 2:
                command, second_command = command.split(",")
            if command == REMOVE_TABBED:
                if "\t" in line:
                    if second_command == TO_GUEST:
                        self.guest.append(line.strip())
                    return ""
            if command == REMOVE_WORD:
                word = trim_item.replace(f"+++{command}+++", "")
                if word in line:
                    trimmed = line.replace(word, "")
                    return " ".join(trimmed.split())
            if command == REMOVE_NUMBER:
                word, num_digits = trim_item.replace(f"+++{command}+++", "").split(
                    "___"
                )
                found = re.match("%s\d{%s,}" % (word, num_digits), line.strip())
                if found:
                    return ""
            if command == CONTAINS_MULTIPLE:
                char, amount = trim_item.replace(f"+++{command}+++", "").split("___")
                count = line.count(char)
                if count >= int(amount):
                    return ""
            if command == TIME_TO_TIME:
                found = (
                    re.search(r"\d+(\:|\.)\d+\s?(–|-)\s?\d+(\:|\.)\d+", line) != None
                )

                if found:
                    return ""
            if command == EXACT:
                word = trim_item.replace(f"+++{command}+++", "")
                if line == word:
                    return ""
            if command == SONGS:
                # Numbered songs
                found = (
                    re.match(
                        r"\d{2}\.\s?([A-Za-z\.\&\']+\s)+\s?(-|-)\s?([A-Za-z\.\&\']+\s)+",
                        line.strip(),
                    )
                    != None
                )
                if found:
                    return ""
            if command == ALL_UPPERCASE:
                found = re.search(r"^\t?[A-ZÁÐÉÍÓÚÝÞÆÖ \.\:\-\+\,]+$", line)
                if found:
                    return ""
            if command == REMOVE_START_TIME:
                found = re.match(r"\t?\d+(\.|\:)\d+(\s|$)", line)
                if found:
                    res = line.replace(found.group(0), "").strip()
                    if second_command == TO_GUEST:
                        self.guest.append(res)
                        return ""
                    return res
            if command == REMOVE_SPEAKER_IDENTIFIER:
                found = re.match(r"\t*[A-ZÁÐÉÍÓÚÝÞÆÖ]+ ?[A-ZÁÐÉÍÓÚÝÞÆÖ]*\:", line)
                if found:
                    res = line.replace(found.group(0), "").strip()
                    if second_command == TO_GUEST:
                        self.guest.append(res)
                        return ""
                    return res

            if command == ONLY_NUMBERS:
                found = re.match(r"^[\d\:\,\.\/\\]+$", line.strip())
                if found:
                    return ""

            if trim_item in line:
                return ""
        return " ".join(line.split())

    def is_trim_command(self, trim: str) -> str:
        """
        Checks if a trim item is a command.
        Returns the command if it is or empty string if not.
        """
        command = re.search(r"\+{3}(.*)\+{3}", trim)
        if command is None:
            return ""
        return command.group(1)

    def read(self, filepath: str) -> list:
        suffix = Path(filepath).suffix
        if ".pptx" in suffix:
            return self.text_from_presentation(filepath)
        if ".pdf" in suffix:
            return self.text_from_pdf(filepath)
        if suffix in [".docx", ".doc", "odt"]:
            return self.text_from_docx(filepath)
        return "Unknown format. Known formats are ['.pdf','.pptx', '.docx', '.doc']"

    def save(self, output_path: str) -> bool:
        # Clean output_path
        output_txt = self.to_txt_path(output_path)

        try:
            with open(output_txt, "w") as f:
                for line in self.text:
                    f.write(line)
                    f.write("\n")

                    # Save guest text if available
            if len(self.guest) > 0:
                output_guest = output_txt.replace(".txt", ".guest")
                with open(output_guest, "w") as f:
                    for line in self.guest:
                        f.write(line)
                        f.write("\n")
            return True

        except IOError:
            print(
                "An error occurred while creating or writing to the file: \
                    {output_txt}"
            )
            return False

    def to_txt_path(self, path: str):
        """Makes a path end with .txt"""
        output = path.split(".")
        output[-1] = "txt"
        return ".".join(output)

    def read_and_trim(self, path: str, trim_list: list) -> list:
        self.guest = []
        text = self.read(path)
        text = self.trim(trim_list)
        # PDFs has odd newline characters, after trimming
        # We need to connect sentences together.
        if ".pdf" in path:
            text = self.connect_sentences(text)
            self.text = text
        return text

    def connect_sentences(self, text: list) -> list:
        # Make it all into one string
        output = " ".join(text)
        # Split on "." to find sentences
        output = re.split(r"\. ", output)
        output = [x + "." for x in output]
        # If last line has a dot, an additional one will be added
        output[-1] = output[-1].replace("..", ".")
        return output


def main():
    file = "document_handler/test_inputs/test_pdf.pdf"
    reader = DocumentReader()
    text = reader.read(file)
    for t in text:
        print(t)
    reader.save("document_handler/test_outputs/saved.txt")


if __name__ == "__main__":
    main()
